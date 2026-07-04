#!/usr/bin/env python3
"""
PPTX Translation Script
Reads a PPTX file and replaces all text content based on a JSON translation mapping,
preserving all formatting: fonts, colors, sizes, positions, images, shapes, connectors, etc.

When translating to English, if the original font is a sans-serif type (common for Chinese fonts),
it will be automatically changed to Arial for better English readability.

Usage:
    python translate_pptx.py <input.pptx> <translations.json> <output.pptx>

The translations JSON maps full paragraph text (all runs concatenated) to translated text:
    {
        "你好": "Hello",
        "谢谢": "Thank you"
    }

For generating the translations JSON from a PPTX, use extract_text.py first.

IMPORTANT — File Organization:
    Both the translations JSON and the output PPTX should be paths inside the current
    working directory (the workspace), NOT inside this skill's own directory. The skill
    detects its own location at runtime via os.path.abspath(__file__) — it works for any
    skill install path, not a hard-coded location. Example:
        python scripts/translate_pptx.py <input.pptx> ./translations.json ./output.pptx
    Never pass a path under the skill's own directory — that location is reserved for
    the skill itself and must remain free of generated artifacts. The script will detect
    such attempts and exit with a clear error message.
"""

import sys
import os
import json
import re
from lxml import etree
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE

# Namespace for DrawingML
NS = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

# Quote normalization mappings
QUOTE_MAP = {
    # Chinese quotes to English quotes
    '‘': "'",  # left single quotation mark
    '’': "'",  # right single quotation mark
    '"': '"',  # left double quotation mark
    '”': '"',  # right double quotation mark
    # Also handle reverse mapping (just in case)
    '"': '"',
    '"': '"',
    ''': "'",
    ''': "'",
}

def normalize_quotes(text):
    """Normalize various quote characters to standard ASCII quotes."""
    if not text:
        return text
    result = []
    for char in text:
        result.append(QUOTE_MAP.get(char, char))
    return ''.join(result)

def normalize_text_for_matching(text):
    """Normalize text for robust matching: quotes, whitespace, etc."""
    if not text:
        return text
    # Normalize quotes
    text = normalize_quotes(text)
    # Normalize whitespace (collapse multiple spaces to single space)
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

# Line spacing: 1.5x ≈ 1.5 (multiple) or ~360 twips (exact)
LINE_SPACING_1_5 = 1.5
LINE_SPACING_1_5_TWIPS = 360  # approximate threshold

# Cache for normalized translation keys (built once on first translate_paragraph call)
_normalized_cache: dict[str, tuple[str, str]] = {}


def contains_words(text):
    """Check if text contains actual words (not just numbers or symbols).
    
    Returns True if the text contains alphabetic characters.
    """
    if not text:
        return False
    return bool(re.search(r'[a-zA-Z]', text))


def limit_font_size(para, max_size=48):
    """Limit font size to max_size for paragraphs that contain words.
    
    Only applies to text that contains actual words (not numbers/symbols only).
    """
    full_text = "".join(run.text for run in para.runs)
    if not contains_words(full_text):
        return
    
    for run in para.runs:
        try:
            if run.font.size and run.font.size > max_size * 12700:  # pts to emu
                run.font.size = max_size * 12700
        except Exception:
            pass


def set_autofit_shape_to_fit_text(text_frame):
    """Set text frame to autofit mode: Resize shape to fit text.
    
    This ensures the shape automatically resizes when translated text is longer
    than the original Chinese text.
    """
    try:
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    except Exception:
        pass


def adjust_line_spacing(para):
    """Adjust paragraph line spacing to 1.5x if original > 1.5x.
    
    Rules:
    - If line_spacing is a float (multiple) and > 1.5 -> set to 1.5
    - If line_spacing is an int (twips) and > ~360 -> set to 1.5
    - If line_spacing <= 1.5 -> keep unchanged
    """
    try:
        ls = para.line_spacing
        if ls is None:
            return
        if isinstance(ls, float) and ls > LINE_SPACING_1_5:
            para.line_spacing = LINE_SPACING_1_5
        elif isinstance(ls, int) and ls > LINE_SPACING_1_5_TWIPS:
            para.line_spacing = LINE_SPACING_1_5
    except Exception:
        pass


def set_font_explicit(run, font_name):
    """Explicitly set both Latin and East Asian fonts in the run's XML.
    This ensures PowerPoint doesn't show 'Use Asian Text Font' for Latin text.
    """
    try:
        rPr = run._r.get_or_add_rPr()
        
        # Set or update a:latin element (Latin text font)
        latin = rPr.find(f'{{{NS["a"]}}}latin')
        if latin is None:
            latin = etree.SubElement(rPr, f'{{{NS["a"]}}}latin')
        latin.set('typeface', font_name)
        
        # Set or update a:ea element (East Asian text font)
        ea = rPr.find(f'{{{NS["a"]}}}ea')
        if ea is None:
            ea = etree.SubElement(rPr, f'{{{NS["a"]}}}ea')
        ea.set('typeface', font_name)
        
        # Also set the main font name attribute on rPr
        rPr.set('typeface', font_name)
        
    except Exception as e:
        # Fallback to simple font.name setting
        try:
            run.font.name = font_name
        except Exception:
            pass

# Common sans-serif fonts that should be changed to Arial when translating to English
# This includes Chinese sans-serif fonts and general sans-serif fonts
SANS_SERIF_FONTS = {
    "微软雅黑", "Microsoft YaHei", "微软雅黑 Light", "Microsoft YaHei Light",
    "Microsoft YaHei UI", "Microsoft YaHei UI Light",
    "黑体", "SimHei", "SimHei Light",
    "Arial", "Arial Narrow", "Arial Black", "Arial Unicode MS",
    "Helvetica", "Helvetica Neue", "Helvetica LT Std",
    "Calibri", "Calibri Light",
    "Verdana", "Verdana Pro",
    "Tahoma", "Tahoma Bold",
    "Trebuchet MS", "Trebuchet MS Bold",
    "Franklin Gothic", "Franklin Gothic Medium",
    "Segoe UI", "Segoe UI Light", "Segoe UI Semibold",
    "Lucida Sans", "Lucida Sans Unicode",
    "Gill Sans", "Gill Sans MT",
    "Optima",
    "Frutiger", "Frutiger LT Std",
    "Myriad Pro", "Myriad Roman",
    "Futura", "Futura Std",
    "Gotham", "Gotham Bold",
    "Open Sans", "OpenSans",
    "Roboto", "Roboto Light",
    "Noto Sans", "Noto Sans CJK SC",
    "Source Sans Pro",
    "Lato", "Lato Light",
    "Montserrat", "Montserrat Light",
    "Deng Xian", "DengXian", "等线",
}

# Common serif fonts that should be changed to Times New Roman when translating to English
SERIF_FONTS = {
    "Times New Roman", "Times", "TimesNR",
    "宋体", "SimSun", "SimSun-ExtB",
    "明体", "MS Mincho", "MS PMincho", "Yu Mincho",
    "Georgia", "Georgia Pro",
    "Garamond", "Garamond Pro", "EB Garamond",
    "Caslon", "Adobe Caslon Pro",
    "Baskerville", "Baskerville Old Face",
    "Palatino", "Palatino Linotype",
    "Bookman", "Bookman Old Style",
    "Cambria",
    "Didot", "Didot LT Std",
    "Bodoni", "Bodoni MT",
    "Century", "Century Schoolbook",
    "Rockwell",
    "Slab", "Slab Serif",
    "Serif", "Transitional", "Old Style",
    "FangSong", "仿宋", "KaiTi", "楷体",
}


def is_sans_serif_font(font_name):
    """Check if a font name is a sans-serif type font."""
    if not font_name:
        return False
    
    # Handle theme fonts (e.g., +mn-ea, +mj-ea, +mn-cs, +mj-cs)
    # These are PowerPoint theme fonts and should be treated as sans-serif for Chinese presentations
    if font_name.startswith('+'):
        # Theme fonts in Chinese presentations are typically sans-serif (e.g., Microsoft YaHei)
        # Check if it's an East Asian font (+ea) or if we can't determine, assume sans-serif
        if '+ea' in font_name or '+cs' in font_name:
            return True
        # For other theme fonts, assume they are sans-serif for safety
        return True
    
    # Check exact match or partial match (case-insensitive)
    font_lower = font_name.lower()
    
    # First check against known sans-serif fonts
    for sans_font in SANS_SERIF_FONTS:
        sans_lower = sans_font.lower()
        if sans_lower in font_lower or font_lower in sans_lower:
            return True
    
    # Additional partial matching for Chinese fonts
    # Check for common Chinese sans-serif font patterns
    if any(pattern in font_lower for pattern in [
        "yahei", "microsoft yahei", "微软雅黑",
        "simhei", "黑体",
        "gothic", "gothic",
        "sans", "sans-serif",
        "arial", "helvetica", "calibri", "verdana", "tahoma",
        "humanist", "geometric", "neo-grotesque",
    ]):
        return True
    
    # Check for common Western sans-serif font patterns
    sans_serif_patterns = [
        "sans", "gothic", "humanist", "geometric", "neo-grotesque",
        "ui", "display", "caption", "text",
    ]
    if any(pattern in font_lower for pattern in sans_serif_patterns):
        return True
    
    return False


def is_serif_font(font_name):
    """Check if a font name is a serif type font."""
    if not font_name:
        return False
    
    # Handle theme fonts
    if font_name.startswith('+'):
        # Theme fonts with serif characteristics
        if any(serif in font_name.lower() for serif in ['serif', 'times', 'song', 'ming']):
            return True
        return False
    
    # Check exact match or partial match (case-insensitive)
    font_lower = font_name.lower()
    
    # First check against known serif fonts
    for serif_font in SERIF_FONTS:
        serif_lower = serif_font.lower()
        if serif_lower in font_lower or font_lower in serif_lower:
            return True
    
    # Additional partial matching for Chinese serif fonts
    if any(pattern in font_lower for pattern in [
        "song", "宋体", "simsun",
        "ming", "明体", "mincho",
        "fangsong", "仿宋", "kaiti", "楷体",
        "times", "georgia", "garamond", "baskerville",
        "palatino", "cambria", "caslon", "bodoni",
    ]):
        return True
    
    # Check for common Western serif font patterns
    serif_patterns = [
        "serif", "times", "georgia", "garamond", "baskerville",
        "palatino", "bookman", "cambria", "didot", "bodoni",
        "century", "rockwell", "slab",
    ]
    if any(pattern in font_lower for pattern in serif_patterns):
        return True
    
    return False


def get_run_font_name(run):
    """Get the font name from a run, checking both the font.name property and XML attributes."""
    try:
        # First try the standard font.name property
        if run.font.name:
            return run.font.name
    except Exception:
        pass
    
    # If font.name is None or empty, check the XML directly
    try:
        rPr = run._r.get_or_add_rPr()
        
        # Check a:latin element (Latin text font)
        latin = rPr.find(f'{{{NS["a"]}}}latin')
        if latin is not None:
            typeface = latin.get('typeface')
            if typeface:
                return typeface
        
        # Check a:ea element (East Asian text font)
        ea = rPr.find(f'{{{NS["a"]}}}ea')
        if ea is not None:
            typeface = ea.get('typeface')
            if typeface:
                return typeface
    except Exception:
        pass
    
    return None


def check_paragraph_theme_font(para):
    """Check if a paragraph uses a theme font (sans-serif).
    Returns True if the paragraph's default font is a sans-serif theme font.
    """
    try:
        # Check the paragraph's _p element for default font settings
        p = para._p
        pPr = p.find(f'{{{NS["a"]}}}pPr')
        if pPr is not None:
            # Check for defRPr (default run properties)
            defRPr = pPr.find(f'{{{NS["a"]}}}defRPr')
            if defRPr is not None:
                latin = defRPr.find(f'{{{NS["a"]}}}latin')
                if latin is not None:
                    typeface = latin.get('typeface')
                    if typeface and is_sans_serif_font(typeface):
                        return True
                
                ea = defRPr.find(f'{{{NS["a"]}}}ea')
                if ea is not None:
                    typeface = ea.get('typeface')
                    if typeface and is_sans_serif_font(typeface):
                        return True
    except Exception:
        pass
    
    return False


def translate_paragraph(para, translations):
    """Translate a paragraph by replacing run text while preserving all formatting.
    If the original font is sans-serif type, change it to Arial for better English readability.

    Returns:
        True if translation was applied
        False if no matching key was found (caller can record as untranslated)
    """
    full_text = "".join(run.text for run in para.runs)
    if not full_text.strip():
        return False

    # Normalize the paragraph text for matching
    normalized_full = normalize_text_for_matching(full_text)

    # Try exact match first
    eng_text: str = ""
    if full_text in translations:
        eng_text = translations[full_text]
    else:
        # Try normalized match (handles quotes, whitespace differences)
        found = False

        # Build normalized translation keys for matching
        if not _normalized_cache:
            for cn_key, en_val in translations.items():
                normalized_key = normalize_text_for_matching(cn_key)
                _normalized_cache[normalized_key] = (cn_key, en_val)

        if normalized_full in _normalized_cache:
            cn_key, eng_text = _normalized_cache[normalized_full]
            found = True

        if not found:
            # Fallback: try whitespace-normalized match only
            normalized_stripped = full_text.strip()
            for cn_key, en_val in translations.items():
                if cn_key.strip() == normalized_stripped:
                    eng_text = en_val
                    found = True
                    break

        if not found:
            return False

    # Check font type and determine target font
    # Check ALL runs in the paragraph for font type
    target_font = None  # None means keep original font

    # Method 1: Check run.font.name using enhanced detection
    for run in para.runs:
        try:
            font_name = get_run_font_name(run)
            if font_name:
                if is_sans_serif_font(font_name):
                    target_font = "Arial"
                    break
                elif is_serif_font(font_name):
                    target_font = "Times New Roman"
                    break
        except Exception:
            continue

    # Method 2: Check paragraph's default font (theme font)
    if target_font is None:
        for run in para.runs:
            try:
                font_name = get_run_font_name(run)
                if font_name is None:  # Theme font
                    if check_paragraph_theme_font(para):
                        # Assume Chinese presentations use sans-serif theme fonts
                        target_font = "Arial"
                        break
            except Exception:
                continue

    # Method 3: Check if no explicit font is set (likely using theme font)
    if target_font is None and para.runs:
        try:
            # If font.name is None for all runs, it's using theme font
            all_none = all(get_run_font_name(run) is None for run in para.runs)
            if all_none:
                # Assume Chinese presentations use sans-serif theme fonts
                target_font = "Arial"
        except Exception:
            pass

    # Put translation in first run, clear the rest
    if para.runs:
        # Apply target font BEFORE setting text to ensure it takes effect
        if target_font:
            set_font_explicit(para.runs[0], target_font)

        para.runs[0].text = eng_text

        # Apply target font to all other runs for consistency
        for run in para.runs[1:]:
            if target_font:
                set_font_explicit(run, target_font)
            run.text = ""

        # Limit font size to 48pt for English text containing words
        if target_font:  # Only when translating Chinese→English
            limit_font_size(para, max_size=48)

    return True


def find_closest_key(text, translation_keys, max_candidates=1):
    """Find translation keys that are most similar to the given text.
    Uses difflib to suggest the closest match — helps the user spot typos
    (e.g., 前提 vs 前要) that prevent exact matching.
    """
    try:
        import difflib
    except ImportError:
        return []

    return difflib.get_close_matches(text, translation_keys, n=max_candidates, cutoff=0.5)


def assert_path_in_workspace(path, label):
    """Refuse to read/write artifacts inside the skill's own directory.

    The skill directory is reserved for the skill itself. All artifacts (JSON, output PPTX)
    must live in the current working directory (workspace) or in a user-chosen path
    OUTSIDE the skill directory. This guard prevents accidental pollution of the skill.
    """
    skill_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    path_abs = os.path.abspath(path)
    skill_dir_norm = os.path.normcase(os.path.normpath(skill_dir))
    path_norm = os.path.normcase(os.path.normpath(path_abs))

    if path_norm == skill_dir_norm or path_norm.startswith(skill_dir_norm + os.sep):
        print(
            f"Error: {label} path is inside the skill directory, which is reserved for the skill itself.\n"
            f"  Path:  {path_abs}\n"
            f"  Skill: {skill_dir}\n"
            f"Please use a path inside your current workspace (e.g. './translations.json' or './output.pptx').",
            file=sys.stderr,
        )
        sys.exit(2)


def process_shape(shape, translations):
    """Recursively process shapes including groups.

    Returns:
        List of (shape, paragraph, full_text) tuples for paragraphs that
        could not be translated. Useful for reporting.
    """
    untranslated = []

    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            full_text = "".join(run.text for run in para.runs)
            if not full_text.strip():
                continue
            if not translate_paragraph(para, translations):
                untranslated.append((shape, para, full_text))
            adjust_line_spacing(para)
        
        # Set autofit: Resize shape to fit text
        set_autofit_shape_to_fit_text(shape.text_frame)

    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    full_text = "".join(run.text for run in para.runs)
                    if not full_text.strip():
                        continue
                    if not translate_paragraph(para, translations):
                        untranslated.append((shape, para, full_text))
                    adjust_line_spacing(para)
                # Set autofit for table cells
                set_autofit_shape_to_fit_text(cell.text_frame)

    # Recursively handle group shapes
    if shape.shape_type == 6:  # GROUP
        for child_shape in shape.shapes:
            untranslated.extend(process_shape(child_shape, translations))

    return untranslated


def main():
    if len(sys.argv) != 4:
        print("Usage: python translate_pptx.py <input.pptx> <translations.json> <output.pptx>")
        sys.exit(1)

    input_path = sys.argv[1]
    translations_path = sys.argv[2]
    output_path = sys.argv[3]

    # Guard: never read/write artifacts inside the skill directory
    assert_path_in_workspace(translations_path, "translations.json")
    assert_path_in_workspace(output_path, "output.pptx")

    # Load translations
    with open(translations_path, "r", encoding="utf-8") as f:
        translations = json.load(f)

    print(f"Loaded {len(translations)} translation entries")

    # Open presentation
    prs = Presentation(input_path)
    print(f"Slides: {len(prs.slides)}, Dimensions: {prs.slide_width}x{prs.slide_height}")

    untranslated = []  # list of (slide_idx, shape_name, full_text, suggestions)

    # Process slides
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            shape_untranslated = process_shape(shape, translations)
            for s, p, text in shape_untranslated:
                shape_name = getattr(s, 'name', '<unnamed>')
                suggestions = find_closest_key(text, list(translations.keys()), max_candidates=1)
                untranslated.append((slide_idx + 1, shape_name, text, suggestions))

    # Also process slide notes
    for slide_idx, slide in enumerate(prs.slides):
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf:
                for para in notes_tf.paragraphs:
                    full_text = "".join(run.text for run in para.runs)
                    if not full_text.strip():
                        continue
                    if not translate_paragraph(para, translations):
                        suggestions = find_closest_key(full_text, list(translations.keys()), max_candidates=1)
                        untranslated.append((slide_idx + 1, "<notes>", full_text, suggestions))
                    adjust_line_spacing(para)
                # Set autofit for notes text frame
                set_autofit_shape_to_fit_text(notes_tf)

    # Save
    prs.save(output_path)
    print(f"Saved: {output_path}")

    if untranslated:
        print(f"\nWarning: {len(untranslated)} paragraph(s) could not be translated (no matching key):")
        for slide_num, shape_name, text, suggestions in untranslated[:10]:
            print(f"  Slide {slide_num} [{shape_name}]: {repr(text[:120])}")
            if suggestions:
                sug = suggestions[0]
                if sug != text:
                    # Show the first divergence so the user can spot typos
                    diff_pos = None
                    for i, (a, b) in enumerate(zip(text, sug)):
                        if a != b:
                            diff_pos = i
                            break
                    if diff_pos is None:
                        diff_pos = min(len(text), len(sug))
                    print(f"      Closest JSON key: {repr(sug[:120])}")
                    print(f"      First difference at char {diff_pos}: PPTX={repr(text[diff_pos:diff_pos+3])} JSON={repr(sug[diff_pos:diff_pos+3])}")
        if len(untranslated) > 10:
            print(f"  ... and {len(untranslated) - 10} more")
        print("\nTip: Translation matching is exact-string (after quote/whitespace normalization).")
        print("     If a paragraph is untranslated, copy the original text from PPTX and use it as the JSON key.")

    print("Done.")


if __name__ == "__main__":
    main()
