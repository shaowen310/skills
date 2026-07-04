#!/usr/bin/env python3
"""
Extract all text from a PPTX file for translation.
Outputs a JSON mapping where keys are original paragraph texts and values are empty strings,
ready to be filled in with translations.

Usage:
    python extract_text.py <input.pptx> <output.json>

This generates a JSON file like:
    {
        "你好": "",
        "谢谢": ""
    }

The translator then fills in the values:
    {
        "你好": "Hello",
        "谢谢": "Thank you"
    }

Then use translate_pptx.py to apply the translations.

IMPORTANT — File Organization:
    The output JSON path must be inside the current working directory (the workspace),
    NOT inside this skill's own directory. The skill detects its own location at runtime
    via os.path.abspath(__file__) — it works for any skill install path, not a hard-coded
    location. Example:
        python scripts/extract_text.py <input.pptx> ./translations_text.json
    Never pass a path under the skill's own directory as the output — that location is
    reserved for the skill itself and must remain free of generated artifacts. The script
    will detect such attempts and exit with a clear error message.
"""

import sys
import os
import json
from pptx import Presentation


def extract_paragraph_text(para, texts_seen):
    """Extract concatenated text of all runs in a paragraph."""
    full_text = "".join(run.text for run in para.runs)
    if full_text.strip() and full_text not in texts_seen:
        texts_seen[full_text] = ""


def process_shape(shape, texts_seen):
    """Recursively process shapes including groups."""
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            extract_paragraph_text(para, texts_seen)

    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    extract_paragraph_text(para, texts_seen)

    if shape.shape_type == 6:  # GROUP
        for child_shape in shape.shapes:
            process_shape(child_shape, texts_seen)


def assert_output_in_workspace(output_path):
    """Refuse to write artifacts into the skill's own directory.

    The skill directory (where this script lives) is reserved for the skill
    itself. Generated artifacts must go to the current working directory (workspace)
    or to a path explicitly chosen by the user that is OUTSIDE the skill directory.
    """
    skill_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    output_abs = os.path.abspath(output_path)
    try:
        # Common path: if the user passed a relative path, this resolves
        # against the current working directory — which is what we want.
        if not os.path.isabs(output_path):
            output_abs = os.path.abspath(output_path)
    except Exception:
        pass

    # Normalize for comparison
    skill_dir_norm = os.path.normcase(os.path.normpath(skill_dir))
    output_norm = os.path.normcase(os.path.normpath(output_abs))

    if output_norm == skill_dir_norm or output_norm.startswith(skill_dir_norm + os.sep):
        print(
            f"Error: Refusing to write generated artifact into the skill directory:\n"
            f"  Output:  {output_abs}\n"
            f"  Skill:   {skill_dir}\n"
            f"Please save the JSON file in your current workspace (e.g. './translations_text.json').",
            file=sys.stderr,
        )
        sys.exit(2)


def main():
    if len(sys.argv) != 3:
        print("Usage: python extract_text.py <input.pptx> <output.json>")
        sys.exit(1)

    input_path = sys.argv[1]
    output_path = sys.argv[2]

    # Guard: never write into the skill directory
    assert_output_in_workspace(output_path)

    prs = Presentation(input_path)
    texts = {}

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            process_shape(shape, texts)

    # Also extract notes
    for slide_idx, slide in enumerate(prs.slides):
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf:
                for para in notes_tf.paragraphs:
                    extract_paragraph_text(para, texts)

    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(texts, f, ensure_ascii=False, indent=2)

    print(f"Extracted {len(texts)} unique text segments from {len(prs.slides)} slides")
    print(f"Saved to: {output_path}")
    print("Fill in the empty string values with translations, then run translate_pptx.py")


if __name__ == "__main__":
    main()
