"""
pptx2md.convert(input_pptx, output_md=None, asset_dir='assets')

Convert a .pptx file to Markdown with images extracted into an assets folder.

Usage:
    from pptx2md import convert
    convert('input.pptx', 'output.md')
    # or in CLI mode
    python -m pptx2md input.pptx output.md

Behavior:
- One slide per section, separated by `---`.
- First short text on a slide (len < 60, no newline) becomes a `##` heading.
- Remaining text is preserved as paragraphs in reading order.
- Pictures (top-level and inside GROUPs) are extracted as `page-XX-img-YY.ext`
  under `asset_dir`. Reading order is reconstructed by (top, left) sort.
- Master / brand pictures are **excluded**: a picture is dropped (not written,
  not referenced, not counted in the index) when ANY of the following hold:
    1. **Full-bleed background** — covers ~95%+ of the slide canvas (typical
       signature of a slide master background).
    2. **Majority-recurring** — identical bytes (sha256) recur on
       `>= max(2, ceil(0.5 * total_slides))` pages (typical signature of a
       master background copied onto every slide).
    3. **Small corner brand asset** — the picture is small relative to the
       slide (area < BRAND_AREA_FRAC * slide area AND width < BRAND_DIM_FRAC
       * slide width AND height < BRAND_DIM_FRAC * slide height) AND sits
       flush in a margin zone (its bounding box starts within BRAND_MARGIN_FRAC
       of a slide edge AND ends within 5% of an opposite slide edge). This
       catches corporate logos and footer/header marks that recur on a small
       subset of section-divider or title pages — pictures that are too
       sparse to be flagged by the majority-recurring rule.
- Returns the absolute path to the produced markdown and a dict of
  {page: [filename, ...]} for downstream tooling. Flowcharts are output as
  cropped rendered images plus a swimlane-structured text fallback.
"""
from __future__ import annotations

import hashlib
import os
import re
import shutil
import subprocess
import sys
import tempfile
from collections import defaultdict
from pathlib import Path
from typing import Any, TypedDict

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# --- Shared data types -------------------------------------------------------
# A slide content item collected from the shape tree. It is heterogeneous — a
# picture, a text block, or a table — so only a subset of fields is present per
# `kind`; all fields are therefore optional. `shape` is the underlying
# python-pptx object, which ships no type stubs, so it stays `Any`.
class SlideItem(TypedDict, total=False):
    kind: str                 # 'text' | 'img'
    text: str
    shape: Any
    hash: str                 # content hash of a picture blob
    ext: str                  # picture file extension
    st: MSO_SHAPE_TYPE        # original shape type
    top: int
    left: int
    width: int
    height: int
    w: int                    # text/table bbox width (text items only)
    h: int                    # text/table bbox height (text items only)


# One extracted picture reference recorded in the image index.
class ImageRef(TypedDict, total=False):
    page: int
    idx: int
    file: str


# A swimlane group produced for diagram slides.
class Lane(TypedDict, total=False):
    name: str
    top: int
    bottom: int
    items: list[SlideItem]


# --- Renderer detection (LibreOffice > PowerPoint COM) -----------------------
_soffice_path: str | None = None


def _find_soffice() -> str | None:
    """Locate the LibreOffice soffice executable.

    Checks (in order): LIBREOFFICE_PATH env var, PATH, then common install
    locations.  Returns the absolute path or None.
    """
    global _soffice_path
    cached = _soffice_path
    if cached is not None:
        return cached

    # 1) Explicit environment variable.
    env = os.environ.get('LIBREOFFICE_PATH')
    if env and os.path.isfile(env):
        _set_soffice(env)
        return _soffice_path

    # 2) PATH.
    exe = shutil.which('soffice')
    if exe:
        _set_soffice(exe)
        return _soffice_path

    # 3) Common install paths.
    common = [
        r'C:\Program Files\LibreOffice\program\soffice.exe',
        r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
        '/usr/lib/libreoffice/program/soffice',
        '/Applications/LibreOffice.app/Contents/MacOS/soffice',
    ]
    for p in common:
        if os.path.isfile(p):
            _set_soffice(p)
            return _soffice_path

    _soffice_path = ''  # sentinel: not found
    return None


def _set_soffice(path: str) -> None:
    global _soffice_path
    _soffice_path = path


def _create_single_slide_pptx(src: str, slide_index: int, out_path: str) -> bool:
    """Create a single-slide .pptx from a specific slide (1-indexed) of the source.

    Saves the slide as the only slide by removing other slide XMLs from the
    OPC zip and adjusting presentation.xml / rels to reference only rId1.
    """
    import zipfile
    from lxml import etree

    NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

    slide_target = f'ppt/slides/slide{slide_index}.xml'
    keep_slides = {slide_target, f'ppt/slides/_rels/slide{slide_index}.xml.rels'}

    with zipfile.ZipFile(src, 'r') as zin:
        names = set(zin.namelist())
        if slide_target not in names:
            return False

        # Read existing rels to find the rId for this slide.
        pres_rels = etree.fromstring(zin.read('ppt/_rels/presentation.xml.rels'))
        slide_rId = None
        for child in pres_rels:
            target = child.get('Target', '').replace('\\', '/')
            # Rels target is relative to ppt/ dir, e.g. "slides/slide6.xml".
            if target.endswith(f'/slide{slide_index}.xml') or target == f'slide{slide_index}.xml':
                slide_rId = child.get('Id')
                break

        if not slide_rId:
            return False

        with zipfile.ZipFile(out_path, 'w', zipfile.ZIP_DEFLATED) as zout:
            for name in names:
                # Skip other slide files.
                if name.startswith('ppt/slides/slide') and name not in keep_slides:
                    continue
                # Skip other slide rels files.
                if name.startswith('ppt/slides/_rels/') and name not in keep_slides:
                    continue

                data = zin.read(name)

                # Patch presentation.xml — keep only the target slide's sldId.
                if name == 'ppt/presentation.xml':
                    root = etree.fromstring(data)
                    sldId_lst = root.find(f'{{{NS_P}}}sldIdLst')
                    if sldId_lst is not None:
                        for child in list(sldId_lst):
                            rId = child.get(f'{{{NS_R}}}id')
                            if rId != slide_rId:
                                sldId_lst.remove(child)
                    data = etree.tostring(root, xml_declaration=True,
                                          encoding='UTF-8', standalone=True)

                zout.writestr(name, data)
    return True


def _render_slides_via_libreoffice(src: str, asset_root: Path,
                                   diagram_pages: set[int]) -> dict[int, str]:
    """Render diagram slides as PNG images using LibreOffice headless mode.

    Each diagram slide is exported as a single-slide PPTX and converted
    individually — avoids ambiguity in LibreOffice's multi-slide PNG naming.
    Returns {page_num: filename} for successfully rendered slides.
    Falls back gracefully (empty dict) on any error.
    """
    soffice = _find_soffice()
    if not soffice or not diagram_pages:
        return {}

    tmpdir: str | None = None
    result: dict[int, str] = {}
    try:
        tmpdir = tempfile.mkdtemp(prefix='pptx2md_')
        for pi in sorted(diagram_pages):
            single_pptx = os.path.join(tmpdir, f'_slide{pi}.pptx')
            if not _create_single_slide_pptx(src, pi, single_pptx):
                continue

            # Convert single-slide PPTX — always produces <stem>.png.
            cmd = [soffice, '--headless', '--convert-to', 'png',
                   '--outdir', tmpdir, single_pptx]
            _ = subprocess.run(cmd, capture_output=True, timeout=60)

            # The output PNG is <stem>.png (single-slide conversion).
            stem = f'_slide{pi}'
            png_path = os.path.join(tmpdir, f'{stem}.png')
            if not os.path.isfile(png_path):
                continue

            target = f'slide-{pi:02d}.png'
            _ = shutil.copy2(png_path, str(asset_root / target))
            result[pi] = target
        return result
    except Exception:
        return {}
    finally:
        if tmpdir and os.path.isdir(tmpdir):
            shutil.rmtree(tmpdir, ignore_errors=True)


def _crop_rendered_image(image_path: str) -> None:
    """Crop a rendered PNG to its content bounding box, removing empty margins.

    Uses Pillow to detect non-white pixels and crops to their bounding box.
    Does nothing if Pillow is unavailable or the image is already tight.
    """
    try:
        from PIL import Image, ImageOps
    except ImportError:
        return
    try:
        img = Image.open(image_path).convert('L')
        # Invert so content (dark) becomes light, white bg becomes dark.
        inv = ImageOps.invert(img)

        def _point_threshold(x: object) -> int:
            # PIL passes an int pixel at runtime; its stub types it as a union
            # that does not support the ">" operator.
            val = x if isinstance(x, (int, float)) else 0
            return 255 if val > 15 else 0

        # Threshold: any pixel not near-white becomes white in the mask.
        mask = inv.point(_point_threshold)
        bbox = mask.getbbox()
        if bbox:
            x1, y1, x2, y2 = bbox
            # Add a 6px margin for aesthetics.
            margin = 6
            x1 = max(0, x1 - margin)
            y1 = max(0, y1 - margin)
            x2 = min(img.width, x2 + margin)
            y2 = min(img.height, y2 + margin)
            # Crop if at least 10px trimmed from any edge.
            if x1 > 10 or y1 > 10 or x2 < img.width - 10 or y2 < img.height - 10:
                colored = Image.open(image_path).crop((x1, y1, x2, y2))
                colored.save(image_path)
    except Exception:
        pass


def _crop_rendered_image_to_bbox(image_path: str,
                                  bbox: tuple[int, int, int, int],
                                  slide_w: int, slide_h: int) -> None:
    """Crop a rendered full-slide PNG down to the flowchart region.

    ``bbox`` is the diagram's bounding box in EMU (see
    :func:`_compute_flowchart_bbox`). The rendered PNG is the whole slide, so
    we map the EMU box into pixel space using the image's actual dimensions
    and crop to it. A small margin is added for aesthetics. Does nothing if
    Pillow is unavailable or the mapping is degenerate.
    """
    try:
        from PIL import Image
    except ImportError:
        return
    try:
        img = Image.open(image_path).convert('RGB')
        W, H = img.size
        if not slide_w or not slide_h:
            return
        sx = W / slide_w
        sy = H / slide_h
        margin = int(0.005 * slide_w)  # ~0.5% of slide width, in EMU
        pl = max(0, int((bbox[0] - margin) * sx))
        pt = max(0, int((bbox[1] - margin) * sy))
        pr = min(W, int((bbox[2] + margin) * sx))
        pb = min(H, int((bbox[3] + margin) * sy))
        if pr - pl < 10 or pb - pt < 10:
            return
        cropped = img.crop((pl, pt, pr, pb))
        cropped.save(image_path)
    except Exception:
        pass


def _shape_top(shape) -> int:
    try:
        return int(shape.top or 0)
    except Exception:
        return 0


def _shape_left(shape) -> int:
    try:
        return int(shape.left or 0)
    except Exception:
        return 0


def _clean_text(text: str) -> str:
    """Remove or replace unreadable/invisible Unicode characters from extracted text.

    PowerPoint text (especially via python-pptx) often contains control characters
    that are invisible in PowerPoint but render as garbled boxes or symbols in
    markdown output. This function sanitizes them so the output is clean text.

    Characters handled:
    - U+000B (vertical tab / manual line break within a paragraph) → space
    - U+000C (form feed) → space
    - U+00A0 (non-breaking space) → regular space
    - C0 control chars (U+0000–U+0008, U+000E–U+001F, U+007F) → stripped
    - C1 control chars (U+0080–U+009F) → stripped
    - Unicode format chars (U+200B–U+200F, U+2028–U+202F, U+FEFF) → stripped
    - Multiple consecutive spaces → collapsed to single space
    """
    if not text:
        return text

    # Replace common PPT control characters with a space.
    text = text.replace('\u000b', ' ')   # vertical tab (manual line break within paragraph)
    text = text.replace('\u000c', ' ')   # form feed
    text = text.replace('\u00a0', ' ')   # non-breaking space

    # Strip other C0 control characters (keep \n, \r, \t which are legit).
    text = re.sub(r'[\x00-\x08\x0e-\x1f\x7f]', '', text)

    # Strip C1 control characters.
    text = re.sub(r'[\x80-\x9f]', '', text)

    # Strip Unicode format / invisible characters.
    text = re.sub(r'[\u200b-\u200f\u2028-\u202f\ufeff]', '', text)

    # Collapse multiple consecutive spaces (but keep intentional line breaks).
    text = re.sub(r'[ \t]+', ' ', text)

    return text.strip()


def _md_relpath(target: Path, base: Path) -> str:
    """Return a relative path from *base* to *target*, using forward slashes.

    Markdown image references require POSIX-style separators (``/``), but
    ``os.path.relpath`` on Windows produces backslashes (``\\``).  This helper
    normalises the result so image links work consistently across platforms.
    """
    return os.path.relpath(target, base).replace('\\', '/')


def _shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ''
    paras = []
    for p in shape.text_frame.paragraphs:
        txt = ''.join(r.text for r in p.runs)
        if not txt and p.text:
            txt = p.text
        paras.append(txt)
    return _clean_text('\n'.join(paras).strip())


def _collect_items(shape, page: int, offset_top: int = 0, offset_left: int = 0) -> list[SlideItem]:
    """Recursively walk shape tree, collecting metadata for text and pictures.

    Pictures are NOT written to disk here. Each picture item carries its
    content hash, extension, geometry, and the shape reference so the caller
    can decide (after global analysis) whether to extract it.
    """
    top = offset_top + _shape_top(shape)
    left = offset_left + _shape_left(shape)
    items = []
    st = shape.shape_type
    if st == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            items.extend(_collect_items(sub, page, top, left))
    elif st == MSO_SHAPE_TYPE.PICTURE:
        try:
            blob = shape.image.blob
            ext = shape.image.ext
            h = hashlib.sha256(blob).hexdigest()
        except Exception:
            return items
        try:
            w = int(shape.width or 0)
            hgt = int(shape.height or 0)
        except Exception:
            w = hgt = 0
        items.append({'kind': 'img', 'hash': h, 'ext': ext,
                      'top': top, 'left': left, 'width': w, 'height': hgt,
                      'shape': shape})
    elif st == MSO_SHAPE_TYPE.TABLE:
        rows_data = []
        num_cols = 0
        for row in shape.table.rows:
            cells = []
            for cell in row.cells:
                # Flatten newlines inside table cells — Markdown table rows are
                # single-line; embedded \n / \r would break the row across
                # multiple lines and render as a dangling orphan line.
                txt = _clean_text(cell.text.replace('\n', ' ').replace('\r', ' '))
                cells.append(txt)
            if len(cells) > num_cols:
                num_cols = len(cells)
            rows_data.append(cells)
        text = _format_table_markdown(rows_data)
        if text:
            items.append({'kind': 'text', 'text': text, 'top': top, 'left': left,
                          'w': int(shape.width or 0), 'h': int(shape.height or 0),
                          'st': MSO_SHAPE_TYPE.TABLE})
    elif shape.has_text_frame:
        text = _shape_text(shape)
        if text:
            items.append({'kind': 'text', 'text': text, 'top': top, 'left': left,
                          'w': int(shape.width or 0), 'h': int(shape.height or 0),
                          'st': st})
    return items


def _format_table_markdown(rows: list[list[str]]) -> str:
    """Format a 2D cell matrix as a GitHub-flavored Markdown table.

    Trailing rows that are entirely empty (PowerPoint often leaves blank rows
    at the bottom of a table, e.g. after the last real entry) are dropped so
    they don't render as dangling ``| | |`` separator lines in the markdown.
    """
    if not rows:
        return ''
    num_cols = max(len(r) for r in rows)
    if num_cols == 0:
        return ''

    # Normalise all rows to the same column count.
    normalised = [r + [''] * (num_cols - len(r)) for r in rows]

    # Drop trailing rows that are entirely empty.
    while normalised and all(c == '' for c in normalised[-1]):
        _ = normalised.pop()
    if not normalised:
        return ''

    # Build the header + separator row + data rows.
    lines = []
    lines.append('| ' + ' | '.join(normalised[0]) + ' |')
    lines.append('| ' + ' | '.join(['---'] * num_cols) + ' |')
    for row in normalised[1:]:
        lines.append('| ' + ' | '.join(row) + ' |')

    return '\n'.join(lines)


# --- Brand / master picture detection tunables -----------------------------
# A picture is treated as a "small corner brand asset" (corporate logo, footer
# mark, watermark) and dropped when it is BOTH small relative to the slide
# AND flush in a margin zone. The defaults are calibrated to catch 16:9 slides
# with corner logos ~0.5"-1" in size without flagging any legitimate
# mid-slide content image. Tune cautiously; tightening risks losing real
# content, loosening risks letting the master logo through.
BRAND_AREA_FRAC = 0.02       # picture area < this fraction of slide area
BRAND_DIM_FRAC = 0.20        # picture width AND height < this fraction of slide
BRAND_MARGIN_FRAC = 0.10     # near-edge: within this fraction of a slide edge
# ---------------------------------------------------------------------------

# Decorative standalone pictures are dropped on slides whose title contains a
# keyword from the curated list in ``references/decorative_keywords.txt``
# (sibling of ``scripts/``). The list is loaded at runtime and is NEVER
# hard-coded here — edit that file (one keyword per line, ``#`` for comments,
# blank lines ignored) to enable/disable the rule per deck. This only affects
# NON-diagram pages and never touches the rendered flowchart image or any body
# paragraph.
DEFAULT_DECORATIVE_KEYWORDS_FILE = (
    Path(__file__).resolve().parent.parent / 'references' / 'decorative_keywords.txt'
)

_decorative_keywords_cache: list[str] | None = None


def _load_decorative_keywords(path: str | Path | None = None) -> list[str]:
    """Load decorative-image title keywords from a text file.

    One keyword per line; lines starting with ``#`` and blank lines are
    ignored. Returns ``[]`` (no exclusions) when the file is missing so the
    script still runs on decks without a configured keyword list.
    """
    p = Path(path) if path else DEFAULT_DECORATIVE_KEYWORDS_FILE
    if not p.exists():
        return []
    kws: list[str] = []
    for line in p.read_text(encoding='utf-8').splitlines():
        line = line.strip()
        if line and not line.startswith('#'):
            kws.append(line)
    return kws


def _title_excluded(title: str, keywords: list[str] | None = None) -> bool:
    """Return True when ``title`` contains a decorative-image keyword.

    When True, a (non-diagram) slide's standalone pictures are skipped so only
    the extracted text is kept. ``keywords`` is resolved once by ``convert()``;
    if omitted, the default references file is loaded (and cached).
    """
    if not title:
        return False
    if keywords is None:
        global _decorative_keywords_cache
        if _decorative_keywords_cache is None:
            _decorative_keywords_cache = _load_decorative_keywords()
        keywords = _decorative_keywords_cache
    return any(kw in title for kw in keywords)


def _is_master_background(item: SlideItem, total_slides: int,
                          hash_pages: dict[str, set[int]], slide_w: int, slide_h: int) -> bool:
    """Return True if the picture is a master/background image to be excluded.

    Three heuristics (any one triggers exclusion):
      1. Full-bleed: covers ~95%+ of the slide canvas.
      2. Repeated: identical image bytes (hash) appear on most slides.
      3. Small corner brand: small picture flush in a slide margin — typically
         a corporate logo or footer mark recurring on a subset of pages
         (e.g. section dividers) that is too sparse for rule 2 to catch.
    """
    try:
        left = item['left']; top = item['top']
        w = item['width'] or 0; hgt = item['height'] or 0
    except Exception:
        return False
    right = left + w; bottom = top + hgt

    # 1) Full-bleed background covering ~95%+ of the canvas.
    if slide_w and slide_h:
        if (abs(left) <= 0.02 * slide_w and abs(top) <= 0.02 * slide_h and
                right >= 0.95 * slide_w and bottom >= 0.95 * slide_h):
            return True

    # 2) Identical image recurring on most slides.
    pages = hash_pages.get(item['hash'], set())
    threshold = max(2, (total_slides + 1) // 2)
    if len(pages) >= threshold:
        return True

    # 3) Small picture flush in a margin zone = corner brand asset.
    if slide_w and slide_h and w > 0 and hgt > 0:
        slide_area = slide_w * slide_h
        pic_area = w * hgt
        small_enough = (
            pic_area < BRAND_AREA_FRAC * slide_area
            and w < BRAND_DIM_FRAC * slide_w
            and hgt < BRAND_DIM_FRAC * slide_h
        )
        if small_enough:
            # Count how many of the picture's four edges sit within
            # BRAND_MARGIN_FRAC of the corresponding slide edge. A small
            # picture anchored to a corner (or running along a side) will
            # have 2 such flush edges; a centered small image will have 0.
            margin = BRAND_MARGIN_FRAC
            flush_edges = 0
            if left <= margin * slide_w:
                flush_edges += 1
            if right >= (1 - margin) * slide_w:
                flush_edges += 1
            if top <= margin * slide_h:
                flush_edges += 1
            if bottom >= (1 - margin) * slide_h:
                flush_edges += 1
            if flush_edges >= 2:
                return True

    return False


# ---------------------------------------------------------------------------


# A slide is treated as a "shape-only diagram" (e.g. SmartArt exported as
# native PowerPoint shapes — no PICTURE, no LINE) and rendered as a
# LibreOffice image when it carries at least this many AUTO_SHAPE descendants.
# Real content slides rarely have more than a handful of decorative shapes,
# while a typical SmartArt diagram has 15-50 nodes. Tuned against the
# real decks where a SmartArt page has 30-40 AUTO_SHAPE nodes but zero
# connectors; ordinary text slides stay well under 5.
SHAPE_DIAGRAM_THRESHOLD = 15


def _count_descendant_shapes(shape, type_filter: set[MSO_SHAPE_TYPE]) -> int:
    """Recursively count shapes of the given types under ``shape``."""
    total = 0
    st = shape.shape_type
    if st in type_filter:
        total += 1
    if st == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            total += _count_descendant_shapes(sub, type_filter)
    return total


def _slide_has_diagram(slide) -> bool:
    """Check if a slide should be rendered as a diagram image.

    Triggers when the slide has either:
      * any LINE/connector shape (the classic flowchart signature), OR
      * a large number of AUTO_SHAPE descendants (covers vector SmartArt
        diagrams that were never converted to PICTURE and contain no
        connectors — e.g. a SmartArt page with no explicit LINE shapes).
    """
    # 1) Connector-based diagram (existing behaviour).
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            return True
    # 2) Shape-density heuristic: SmartArt / vector diagrams usually have
    #    many AUTO_SHAPE nodes even without explicit LINE connectors.
    auto_count = 0
    for shape in slide.shapes:
        auto_count += _count_descendant_shapes(
            shape, {MSO_SHAPE_TYPE.AUTO_SHAPE})
    return auto_count >= SHAPE_DIAGRAM_THRESHOLD


def _compute_flowchart_bbox(slide) -> tuple[int, int, int, int] | None:
    """Bounding box (left, top, right, bottom) in EMU of the flowchart itself.

    Built from LINE connectors **and** the AUTO_SHAPE nodes those connectors
    actually touch (recursing through GROUPs with accumulated offsets). This
    deliberately excludes standalone label / container boxes — e.g. the
    "业务场景" / "流程说明" caption boxes on the left of a split slide, or a
    big background frame — so the explanatory copy is treated as living
    *outside* the diagram. The bbox is used to (a) decide which text items are
    explanatory copy versus node labels that live *inside* the diagram, and
    (b) crop the rendered slide image down to just the diagram region so the
    descriptive text on the left is not duplicated in the picture.

    Returns ``None`` when the slide carries no connectors or connector-touched
    auto-shapes.
    """
    # [left, top, right, bottom]; empty until a connector/node is seen.
    bbox: list[int] = []
    # Tolerance (EMU) for "a connector touches this node": connectors usually
    # meet a node edge exactly, but a small gap is common.
    tol = 100_000

    def intersects(a, b) -> bool:
        return not (a[2] < b[0] - tol or a[0] > b[2] + tol or
                    a[3] < b[1] - tol or a[1] > b[3] + tol)

    def upd(l: int, t: int, r: int, b: int) -> None:
        if not bbox:
            bbox[:] = [l, t, r, b]
        else:
            bbox[0] = min(bbox[0], l)
            bbox[1] = min(bbox[1], t)
            bbox[2] = max(bbox[2], r)
            bbox[3] = max(bbox[3], b)

    # Pass 1: collect every connector bbox (so we can test node adjacency).
    conn: list[tuple[int, int, int, int]] = []

    def collect_conn(shape, off_l: int, off_t: int) -> None:
        st = shape.shape_type
        if st == MSO_SHAPE_TYPE.GROUP:
            for sub in shape.shapes:
                collect_conn(sub, off_l + int(shape.left or 0),
                             off_t + int(shape.top or 0))
            return
        if st == MSO_SHAPE_TYPE.LINE:
            l = off_l + int(shape.left or 0)
            t = off_t + int(shape.top or 0)
            r = l + int(shape.width or 0)
            b = t + int(shape.height or 0)
            conn.append((l, t, r, b))

    for shape in slide.shapes:
        collect_conn(shape, 0, 0)

    # Pass 2: expand the bbox with connectors and every node a connector
    # touches (a node whose bbox overlaps a connector bbox within `tol`).
    def collect_nodes(shape, off_l: int, off_t: int) -> None:
        st = shape.shape_type
        if st == MSO_SHAPE_TYPE.GROUP:
            gl = off_l + int(shape.left or 0)
            gt = off_t + int(shape.top or 0)
            gr = gl + int(shape.width or 0)
            gb = gt + int(shape.height or 0)
            # A grouped sub-diagram (e.g. a vertical stage-label band or a
            # swimlane container) often carries no connectors of its own yet
            # still belongs to the flow. If the group's own bbox is touched
            # by a connector, grow the diagram bbox to include it so its
            # labels are not cropped away.
            if conn and any(intersects((gl, gt, gr, gb), c) for c in conn):
                upd(gl, gt, gr, gb)
            for sub in shape.shapes:
                collect_nodes(sub, gl, gt)
            return
        if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
            l = off_l + int(shape.left or 0)
            t = off_t + int(shape.top or 0)
            r = l + int(shape.width or 0)
            b = t + int(shape.height or 0)
            if conn and any(intersects((l, t, r, b), c) for c in conn):
                upd(l, t, r, b)

    for c in conn:
        upd(*c)
    for shape in slide.shapes:
        collect_nodes(shape, 0, 0)

    if not bbox:
        return None
    return (bbox[0], bbox[1], bbox[2], bbox[3])


def _extract_lanes(slide, items: list[SlideItem]) -> list[Lane]:
    """Group items on a diagram slide into swimlanes using GROUP shapes as boundaries.

    Returns a list of lane dicts: {'name': str, 'items': [item, ...]}.
    The first entry is the 'role bar' (items above any swimlane).
    """
    # Collect GROUP positions (swimlane headers).
    lanes_raw: list[Lane] = []
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.GROUP:
            continue
        # Find the first text inside the group — that's the lane name.
        lane_name = ''
        for sub in shape.shapes:
            if sub.has_text_frame:
                for p in sub.text_frame.paragraphs:
                    t = ''.join(r.text for r in p.runs)
                    if t.strip():
                        lane_name = t.strip()
                        break
                if lane_name:
                    break
        if lane_name:
            top = int(shape.top or 0)
            bottom = top + int(shape.height or 0)
            lanes_raw.append({'name': lane_name, 'top': top, 'bottom': bottom})

    # Sort lanes by vertical position.
    lanes_raw.sort(key=lambda x: x['top'])

    # Build lane groups — items belong to the last lane whose top is <= item top.
    grouped: list[Lane] = []
    if lanes_raw:
        # Items above the first lane form the "角色" bar.
        grouped.append({'name': '**角色**', 'items': []})  # role bar
        for lane in lanes_raw:
            lane['items'] = []
            grouped.append(lane)
    else:
        # No lanes detected — all items stay together.
        return [{'name': '', 'items': items}]

    for it in items:
        it_top = it.get('top', 0)
        # Determine which lane this item belongs to.
        assigned = False
        for lane in reversed(lanes_raw):
            if it_top >= lane['top']:
                lane['items'].append(it)
                assigned = True
                break
        if not assigned:
            # Above all lanes → role bar.
            grouped[0]['items'].append(it)

    # Remove empty groups.
    return [g for g in grouped if g['items']]


def _format_diagram_slide(slide, items: list[SlideItem]) -> tuple[list[str], bool]:
    """Format a diagram slide as structured swimlane markdown.

    Returns (md_lines, True) on success, (md_lines, False) if fallback to flat text.
    """
    lanes = _extract_lanes(slide, items)
    if not lanes:
        return ([], False)

    # Collect lane names to filter duplicates from items.
    lane_names = {l['name'].lstrip('*').strip() for l in lanes if l.get('name')}

    # Sort items by visual flow within each lane.
    # Discretize top (group within ~0.5mm / 50000 EMU) so near-identical rows
    # sort by left alone, preserving the true horizontal reading order.
    ROW_TOLERANCE = 50000
    lines: list[str] = []
    title_done = False
    for lane in lanes:
        lane['items'].sort(key=lambda x: (x.get('top', 0) // ROW_TOLERANCE, x.get('left', 0)))
        lane_texts: list[str] = []
        for it in lane['items']:
            t = it.get('text', '').strip()
            if not t:
                continue
            # Treat this as the slide title (first short text overall).
            if not title_done and len(t) < 60 and '\n' not in t:
                lines.append(f'## {t}')
                title_done = True
                continue
            # Skip items that are just the lane header label itself.
            if t in lane_names:
                continue
            lane_texts.append(t)

        if lane_texts:
            label = lane.get('name', '')
            if label:
                lines.append(f'### {label}')
            # Join items with → to show flow direction.
            lines.append(' → '.join(lane_texts))
            lines.append('')

    return (lines, True) if lines else ([], False)


def convert(input_pptx: str, output_md: str | None = None,
            asset_dir: str = 'assets',
            decorative_keywords: list[str] | str | None = None
            ) -> tuple[str, dict[int, list[str]]]:
    """
    Convert a .pptx to Markdown.

    Parameters
    ----------
    input_pptx : str
        Absolute or relative path to the source .pptx file.
    output_md : str, optional
        Absolute or relative path to the output .md file. If None, the md
        filename is derived from the source pptx (same stem).
    asset_dir : str
        Directory (relative to the md file or absolute) to hold extracted
        images. Created if missing; old contents are removed.
    decorative_keywords : list[str] | str | None
        Controls dropping of standalone decorative pictures on non-diagram
        slides by title match (see "Decorative standalone-picture exclusion").
        * ``None`` (default) → load keywords from the skill's
          ``references/decorative_keywords.txt``.
        * a ``list``/``tuple`` of strings → use these keywords directly.
        * a ``str``/``Path`` → path to a keyword text file to load.
        Pass an empty list to disable the rule entirely.

    Returns
    -------
    (md_path, image_index_by_page) : tuple
        `image_index_by_page` is a dict {page: [filename, ...]}.
    """
    if decorative_keywords is None:
        decorative_keywords = _load_decorative_keywords()
    elif isinstance(decorative_keywords, (str, Path)):
        decorative_keywords = _load_decorative_keywords(decorative_keywords)
    else:
        decorative_keywords = list(decorative_keywords)

    src = Path(input_pptx).resolve()
    if not src.exists():
        raise FileNotFoundError(src)

    md_path = Path(output_md) if output_md else src.with_suffix('.md')
    md_path = md_path.resolve()

    asset_root = Path(asset_dir)
    if not asset_root.is_absolute():
        asset_root = (md_path.parent / asset_root).resolve()
    if asset_root.exists():
        for f in asset_root.iterdir():
            if f.is_file():
                f.unlink()
    asset_root.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(src))
    image_index: list[ImageRef] = []
    slide_markdowns: list[str] = []

    slide_w = prs.slide_width or 0
    slide_h = prs.slide_height or 0

    # Phase 1: collect items per slide and every picture's content hash.
    slide_items: dict[int, list[SlideItem]] = {}
    hash_pages: dict[str, set[int]] = defaultdict(set)
    diagram_pages: set[int] = set()
    flowchart_bbox: dict[int, tuple[int, int, int, int] | None] = {}
    for pi, slide in enumerate(prs.slides, 1):
        items: list[SlideItem] = []
        for shape in slide.shapes:
            items.extend(_collect_items(shape, pi))
        slide_items[pi] = items
        for it in items:
            if it['kind'] == 'img':
                hash_pages[it['hash']].add(pi)
        # Detect diagram slides (those with connector lines).
        if _slide_has_diagram(slide):
            diagram_pages.add(pi)
            flowchart_bbox[pi] = _compute_flowchart_bbox(slide)
    total_slides = len(prs.slides)

    # Build a lookup so we can access the original slide objects in Phase 2.
    slides_lookup: dict[int, Any] = {}
    for pi, slide in enumerate(prs.slides, 1):
        slides_lookup[pi] = slide

    # Try rendering diagram slides as images (LibreOffice).
    rendered: dict[int, str] = _render_slides_via_libreoffice(
        str(src), asset_root, diagram_pages)

    # Crop rendered images to the flowchart region (or content bbox as a
    # fallback when no diagram geometry was found).
    if rendered:
        for pi, fname in list(rendered.items()):
            path = str(asset_root / fname)
            fb = flowchart_bbox.get(pi)
            if fb:
                _crop_rendered_image_to_bbox(path, fb, slide_w, slide_h)
            else:
                _crop_rendered_image(path)

    # Phase 2: drop master/background pictures, keep only real content.
    for pi in sorted(slide_items):
        items = slide_items[pi]
        items.sort(key=lambda x: (x['top'], x['left']))

        md_lines: list[str] = []
        title_done = False
        img_no = 0

        if pi in diagram_pages:
            rendered_file = rendered.get(pi)
            if rendered_file:
                # Rendered image available — use it.
                slide_obj = slides_lookup.get(pi)
                title_text = ''
                if slide_obj:
                    for it in items:
                        if it['kind'] == 'text':
                            t = it.get('text', '').strip()
                            if t and len(t) < 60 and '\n' not in t:
                                title_text = t
                                break
                if title_text:
                    md_lines.append(f'## {title_text}')
                # Split the slide into explanatory text (which lives OUTSIDE
                # the diagram — typically the left-hand "业务场景" / "流程说明"
                # columns and any captions) and the flowchart itself. The
                # explanatory text is emitted as paragraphs FIRST, then the
                # cropped flowchart image is appended, so the descriptive copy
                # is never interleaved with — or buried under — the diagram.
                # Flowchart node labels are rendered inside the image and are
                # intentionally NOT dumped as loose text.
                fb = flowchart_bbox.get(pi)
                # Clearance (EMU) used to decide a text block is a side column
                # rather than text that lives inside the diagram.
                gap = int(0.08 * slide_w) if slide_w else 0

                def _overlaps_diagram(it: SlideItem) -> bool:
                    if fb is None:
                        return False
                    il = it.get('left', 0)
                    it_ = it.get('top', 0)
                    ir = il + it.get('w', 0)
                    ib = it_ + it.get('h', 0)
                    return not (ir < fb[0] or il > fb[2] or ib < fb[1] or it_ > fb[3])

                # A "split" slide carries a left caption column separated from
                # the flowchart by a real gap. On those slides the explanatory
                # text is pulled out and emitted BEFORE the diagram. On a pure
                # flowchart slide (the whole slide IS the diagram) nothing
                # extra is emitted — the rendered image already carries the
                # swimlane headers, node labels and all other content.
                split = False
                if fb is not None:
                    text_rights = [
                        it.get('left', 0) + it.get('w', 0)
                        for it in items
                        if it['kind'] == 'text'
                        and it.get('text', '').strip()
                        and it.get('text', '').strip() != title_text.strip()
                    ]
                    if text_rights and min(text_rights) < fb[0] - gap:
                        split = True

                descriptive: list[SlideItem] = []
                for it in items:
                    if it['kind'] != 'text':
                        continue
                    t = it.get('text', '').strip()
                    if not t or t == title_text.strip():
                        continue
                    if not split:
                        continue
                    # `split` is only True when `fb is not None` (see above),
                    # so the diagram bbox is safe to dereference here.
                    assert fb is not None
                    # Explanatory copy = free text boxes (captions /
                    # descriptions) plus any shape text that sits clearly to
                    # the LEFT of the flowchart (a side column). Everything
                    # else lives inside the diagram and is left to the image.
                    if (it.get('st') == MSO_SHAPE_TYPE.TEXT_BOX
                            or it.get('left', 0) < fb[0] - gap):
                        descriptive.append(it)

                # Pair short "section labels" with the descriptive paragraphs
                # that follow them (e.g. 业务场景 + its paragraph, then
                # 流程说明 + its paragraph) so the prose reads naturally.
                labels = [d for d in descriptive if len(d['text'].strip()) <= 6]
                paras = [d for d in descriptive if len(d['text'].strip()) > 6]
                labels.sort(key=lambda x: (x.get('top', 0), x.get('left', 0)))
                paras.sort(key=lambda x: (x.get('top', 0), x.get('left', 0)))
                if labels and paras:
                    ordered_descriptive: list[SlideItem] = []
                    for i in range(max(len(labels), len(paras))):
                        if i < len(labels):
                            ordered_descriptive.append(labels[i])
                        if i < len(paras):
                            ordered_descriptive.append(paras[i])
                else:
                    ordered_descriptive = sorted(
                        descriptive,
                        key=lambda x: (x.get('top', 0), x.get('left', 0)))

                for d in ordered_descriptive:
                    md_lines.append('')
                    md_lines.append(d['text'].strip())

                # The flowchart itself is left to the cropped rendered image;
                # its nodes are not reliably extractable as text, so we attach
                # the picture rather than attempt a (high-error) swimlane parse.
                rel = _md_relpath(asset_root / rendered_file, md_path.parent)
                md_lines.append('')
                md_lines.append(f'![Image]({rel})')

                # Pictures that live inside the diagram are already part of the
                # rendered image; only standalone pictures outside it remain.
                pic_items = [it for it in items if it['kind'] == 'img']
                for it in pic_items:
                    if _is_master_background(it, total_slides, hash_pages,
                                             slide_w, slide_h):
                        continue
                    if _overlaps_diagram(it):
                        continue
                    # Decorative standalone pictures on keyword-matched pages
                    # are dropped — text suffices (keywords from references file).
                    if _title_excluded(title_text, decorative_keywords):
                        continue
                    img_no += 1
                    ext = it['ext']
                    new_name = f'page-{pi:02d}-img-{img_no:02d}.{ext}'
                    out_path = asset_root / new_name
                    with open(out_path, 'wb') as f:
                        _ = f.write(it['shape'].image.blob)
                    image_index.append({'page': pi, 'idx': img_no, 'file': new_name})
                    rel2 = _md_relpath(asset_root / new_name, md_path.parent)
                    md_lines.append('')
                    md_lines.append(f'![Image]({rel2})')
                image_index.append({'page': pi, 'idx': 0, 'file': rendered_file})
                slide_markdowns.append(
                    f'<!-- Slide {pi} -->\n\n' + '\n'.join(md_lines).rstrip() + '\n')
                continue
            else:
                # Rendering unavailable — use structured swimlane text.
                slide_obj = slides_lookup.get(pi)
                if slide_obj:
                    diag_lines, used_lanes = _format_diagram_slide(slide_obj, items)
                    if used_lanes:
                        pic_items = [it for it in items if it['kind'] == 'img']
                        pic_lines: list[str] = []
                        for it in pic_items:
                            if _is_master_background(it, total_slides, hash_pages, slide_w, slide_h):
                                continue
                            img_no += 1
                            ext = it['ext']
                            new_name = f'page-{pi:02d}-img-{img_no:02d}.{ext}'
                            out_path = asset_root / new_name
                            with open(out_path, 'wb') as f:
                                _ = f.write(it['shape'].image.blob)
                            image_index.append({'page': pi, 'idx': img_no, 'file': new_name})
                            rel = _md_relpath(asset_root / new_name, md_path.parent)
                            pic_lines.append(f'![Image]({rel})')
                        md_lines = diag_lines + pic_lines
                        slide_markdowns.append(
                            f'<!-- Slide {pi} -->\n\n' + '\n'.join(md_lines).rstrip() + '\n')
                        continue

        slide_title = ''
        for it in items:
            if it['kind'] == 'text':
                t = it['text']
                if not title_done and len(t) < 60 and '\n' not in t:
                    md_lines.append(f'## {t}')
                    title_done = True
                    slide_title = t
                    continue
                title_done = True
                md_lines.append(t)
                md_lines.append('')
            else:
                # Skip master/background pictures (full-bleed or repeated).
                if _is_master_background(it, total_slides, hash_pages, slide_w, slide_h):
                    continue
                # Drop standalone decorative pictures on keyword-matched pages
                # — the extracted text is enough (keywords from references file).
                if _title_excluded(slide_title, decorative_keywords):
                    continue
                img_no += 1
                ext = it['ext']
                new_name = f'page-{pi:02d}-img-{img_no:02d}.{ext}'
                out_path = asset_root / new_name
                with open(out_path, 'wb') as f:
                    _ = f.write(it['shape'].image.blob)
                image_index.append({'page': pi, 'idx': img_no, 'file': new_name})
                # Use a relative path so the md is portable
                rel = _md_relpath(asset_root / new_name, md_path.parent)
                md_lines.append(f'![Image]({rel})')
                md_lines.append('')

        if not md_lines:
            md_lines = [f'<!-- Blank slide {pi} -->']
        slide_markdowns.append(f'<!-- Slide {pi} -->\n\n' + '\n'.join(md_lines).rstrip() + '\n')

    with open(md_path, 'w', encoding='utf-8') as f:
        _ = f.write(f'# {src.stem}\n\n')
        _ = f.write('\n---\n\n'.join(slide_markdowns))

    by_page: dict[int, list[str]] = defaultdict(list)
    for entry in image_index:
        by_page[entry['page']].append(entry['file'])
    return str(md_path), dict(by_page)


def main(argv: list[str]) -> int:
    if len(argv) < 2:
        print('usage: python -m pptx2md <input.pptx> [output.md] [--keywords-file PATH]', file=sys.stderr)
        return 1
    keywords_file = None
    args = list(argv[1:])
    if '--keywords-file' in args:
        i = args.index('--keywords-file')
        if i + 1 >= len(args):
            print('--keywords-file requires a PATH', file=sys.stderr)
            return 1
        keywords_file = args[i + 1]
        del args[i:i + 2]
    md_path, idx = convert(args[0], args[1] if len(args) >= 2 else None,
                           decorative_keywords=keywords_file)
    total = sum(len(v) for v in idx.values())
    print(f'md: {md_path}')
    print(f'images: {total} across {len(idx)} pages')
    return 0


if __name__ == '__main__':
    sys.exit(main(sys.argv))
